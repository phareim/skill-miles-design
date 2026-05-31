#!/usr/bin/env python3
"""
Generate examples/slides/miles-templates.pptx — a set of Miles slide templates
in real PowerPoint, mirroring the HTML deck (deck.html) and the patterns in
../../presentation-patterns.md.

This is a *system demo / starter*, NOT a replacement for the official
`Miles template 2026.pptx`, which remains the source of truth (it has locked
master layouts, real photography, the full 43-layout set). This file proves the
tokens + assets map cleanly into .pptx and gives a lightweight starting point.

Fonts follow the official template's choice: Manrope for headings (Gelica's EULA
forbids .pptx embedding), DM Sans for body. Those fonts must be installed on the
machine that opens the file; they are not embedded here.

Run:
    pip install python-pptx
    python build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import tufte_pptx as tv  # native Tufte-styled chart builders

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.normpath(os.path.join(HERE, "..", "..", "assets"))

# ---- Brand tokens ----
KREM = RGBColor(0xFB, 0xF0, 0xE5)
KREM_DEEP = RGBColor(0xF3, 0xE3, 0xD0)
BURG = RGBColor(0x45, 0x0D, 0x21)
ROD = RGBColor(0xFF, 0x30, 0x3B)
HEAD = "Manrope"          # PPTX heading font (sanctioned Gelica fallback)
BODY = "DM Sans"          # body font

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
MX = Inches(0.85)         # side margin
MY = Inches(0.55)         # top/bottom margin

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=KREM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, left, top, width, height, runs, size=18, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.2, italic=False,
         letter=None):
    """runs: str, or list of (text, RGBColor, bold) for mixed-colour lines."""
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, None, bold)]
    for t, color, b in runs:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.name = font
        r.font.bold = bool(b)
        r.font.italic = italic
        r.font.color.rgb = color or BURG
        if letter is not None:
            _spacing(r, letter)
    return tb


def _spacing(run, pts):
    run.font._rPr.set('spc', str(int(pts * 100)))


def rrect(s, left, top, width, height, fill=None, lineclr=None, linew=1.5,
          radius=0.5):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if lineclr is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = lineclr; sp.line.width = Pt(linew)
    sp.shadow.inherit = False
    return sp


def pill(s, left, top, text_, fill, fg, size=12, width=None, pad=Inches(0.18)):
    w = width or (Inches(0.10) * len(text_) + pad * 2)
    h = Inches(0.42)
    sp = rrect(s, left, top, w, h, fill=fill, lineclr=(BURG if fill is None else None),
               radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text_
    r.font.size = Pt(size); r.font.name = BODY; r.font.bold = True
    r.font.color.rgb = fg
    return sp


def oval(s, left, top, d, fill=None, lineclr=BURG, linew=1.5):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = lineclr; sp.line.width = Pt(linew)
    sp.shadow.inherit = False
    return sp


def img(s, path, left, top, width=None, height=None):
    return s.shapes.add_picture(os.path.join(ASSETS, path), left, top,
                                width=width, height=height)


def chrome(s, dark=False):
    """Hamburger top-left, Miles wordmark top-right."""
    clr = KREM if dark else BURG
    for k in range(3):
        ln = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, MY + Inches(0.06 * k),
                                Inches(0.30), Pt(2.2))
        ln.fill.solid(); ln.fill.fore_color.rgb = clr
        ln.line.fill.background(); ln.shadow.inherit = False
    logo = "logos/miles-logo-cream.png" if dark else "logos/miles-logo-red.png"
    w = Inches(1.15)
    img(s, logo, EMU_W - MX - w, MY - Inches(0.03), width=w)


def kicker(s, left, top, txt, dark=False):
    text(s, left, top, Inches(6), Inches(0.5),
         [(txt, KREM if dark else ROD, False)], size=15, font=HEAD, italic=True)


# ============================== SLIDES ==============================

# 1 · Cover cream
s = slide()
chrome(s)
text(s, MX, Inches(4.7), Inches(6), Inches(0.4), [("Oslo · 2026", ROD, False)], size=13, font=BODY)
img(s, "logos/miles-logo-red.png", MX, Inches(5.05), width=Inches(5.0))
text(s, MX, Inches(6.55), Inches(6), Inches(0.5), "Slide-maler i Miles' designsystem.", size=18, font=HEAD)

# 2 · Cover red flood
s = slide(ROD)
chrome(s, dark=True)
text(s, MX, Inches(4.7), Inches(6), Inches(0.4), [("Bergen · 2026", KREM, False)], size=13)
img(s, "logos/miles-logo-cream.png", MX, Inches(5.05), width=Inches(5.0))

# 3 · Cover burgundy flood
s = slide(BURG)
chrome(s, dark=True)
text(s, MX, Inches(4.7), Inches(6), Inches(0.4), [("Trondheim · 2026", KREM, False)], size=13)
img(s, "logos/miles-logo-cream.png", MX, Inches(5.05), width=Inches(5.0))

# 4 · Illustration cover
s = slide()
chrome(s)
text(s, MX, Inches(2.5), Inches(6.3), Inches(2.0),
     [("Vi elsker å løse\n", BURG, False), ("utfordringer", ROD, False), ("!", BURG, False)],
     size=44, font=HEAD, bold=True, line=1.0)
text(s, MX, Inches(4.7), Inches(5.0), Inches(1.0),
     "Et verdidrevet konsulentselskap med varme for hverandre.", size=18, font=HEAD)
img(s, "illustrations/catching-ideas-net.png", Inches(8.0), Inches(1.7), height=Inches(4.4))

# 5 · TOC
s = slide()
chrome(s)
kicker(s, MX, MY + Inches(0.7), "Innholdsoversikt")
text(s, MX, MY + Inches(1.1), Inches(9), Inches(0.8), "Hva vi skal innom", size=30, font=HEAD, bold=True)
cards = [("Del 1: Bli kjent", ["Hvem vi er", "Våre verdier"]),
         ("Del 2: Tjenester", ["Data og AI", "Sky og plattform"]),
         ("Del 3: Samarbeid", ["Teamet", "Veien videre"])]
cw, gap = Inches(3.6), Inches(0.35)
x0 = MX
for i, (title, items) in enumerate(cards):
    cx = x0 + i * (cw + gap)
    rrect(s, cx, Inches(3.0), cw, Inches(3.3), fill=None, lineclr=BURG, radius=0.06)
    pill(s, cx + Inches(0.3), Inches(3.3), title, fill=BURG, fg=KREM, size=13)
    for j, it in enumerate(items):
        pill(s, cx + Inches(0.3), Inches(4.05) + j * Inches(0.6), it, fill=None, fg=BURG, size=12)

# 6 · Chapter divider red, title in frame
s = slide(ROD)
chrome(s, dark=True)
fr = rrect(s, Inches(2.6), Inches(2.5), Inches(8.1), Inches(2.5), fill=None, lineclr=KREM, radius=0.08)
text(s, Inches(3.0), Inches(2.85), Inches(7.3), Inches(0.4), [("Kapittel 02", KREM, False)], size=14, align=PP_ALIGN.CENTER)
text(s, Inches(3.0), Inches(3.35), Inches(7.3), Inches(1.4), [("Tjenesteområdene våre", KREM, False)], size=40, font=HEAD, bold=True, align=PP_ALIGN.CENTER, line=1.0)

# 7 · Standard content
s = slide()
chrome(s)
kicker(s, MX, Inches(2.2), "Slik jobber vi")
text(s, MX, Inches(2.65), Inches(6.0), Inches(1.4),
     [("Smarte løsninger i ", BURG, False), ("samarbeid", ROD, False)], size=34, font=HEAD, bold=True, line=1.05)
text(s, MX, Inches(4.0), Inches(5.6), Inches(2.0),
     "Vi møter mennesker med faglig autoritet og varme. Vi lytter, forstår problemet, og bygger noe robust og varig sammen med dere.",
     size=18, font=BODY, line=1.5)
img(s, "illustrations/person-laptop-orbit.png", Inches(8.1), Inches(1.9), height=Inches(4.0))

# 8 · Service-area intro
s = slide()
chrome(s)
text(s, MX, Inches(2.2), Inches(6), Inches(0.4), [("TJENESTEOMRÅDE", ROD, False)], size=12, font=BODY, bold=True, letter=1.5)
text(s, MX, Inches(2.65), Inches(6), Inches(1.2), "Data og AI", size=48, font=HEAD, bold=True)
text(s, MX, Inches(6.2), Inches(6.5), Inches(0.8),
     "Maskinlæring · dataplattformer · MLOps · generativ AI · beslutningsstøtte · datakvalitet",
     size=14, font=BODY)
img(s, "service-icons/data-and-ai-circle.png", Inches(8.7), Inches(2.0), height=Inches(3.4))

# 9 · All service areas
s = slide()
chrome(s)
kicker(s, Inches(0), Inches(2.4), "Våre tjenesteområder")
# centre the kicker
text(s, Inches(0), Inches(2.4), EMU_W, Inches(0.5), [("Våre tjenesteområder", ROD, False)],
     size=18, font=HEAD, italic=True, align=PP_ALIGN.CENTER)
svcs = [("strategic-it-circle.png", "Strategisk IT"),
        ("cloud-circle.png", "Teknologi, sky og plattform"),
        ("data-and-ai-circle.png", "Data og AI"),
        ("ux-and-innovation-circle.png", "Brukeropplevelse og innovasjon"),
        ("transformation-and-people-circle.png", "Transformasjon og mennesker")]
n = len(svcs); d = Inches(1.5); gapx = (EMU_W - 2 * MX - n * d) / (n - 1)
for i, (f, label) in enumerate(svcs):
    cx = MX + i * (d + gapx)
    img(s, "service-icons/" + f, cx, Inches(3.2), width=d)
    text(s, cx - Inches(0.25), Inches(4.85), d + Inches(0.5), Inches(0.8),
         [(label, BURG, False)], size=13, font=BODY, align=PP_ALIGN.CENTER, line=1.15)

# 10 · Red split panel
s = slide()
half = EMU_W / 2
left_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, half, EMU_H)
left_rect.fill.solid(); left_rect.fill.fore_color.rgb = ROD; left_rect.line.fill.background(); left_rect.shadow.inherit = False
right_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, half, 0, half, EMU_H)
right_rect.fill.solid(); right_rect.fill.fore_color.rgb = KREM; right_rect.line.fill.background(); right_rect.shadow.inherit = False
text(s, MX, Inches(2.6), half - MX - Inches(0.4), Inches(0.4), [("Mye er nytt", KREM, False)], size=15, font=HEAD, italic=True)
text(s, MX, Inches(3.05), half - MX - Inches(0.4), Inches(1.2), [("En ny verktøykasse", KREM, False)], size=32, font=HEAD, bold=True)
text(s, MX, Inches(4.1), half - MX - Inches(0.4), Inches(1.0), [("Nye plattformer, ny AI, nye muligheter.", KREM, False)], size=17, font=BODY)
text(s, half + Inches(0.4), Inches(2.6), half - MX - Inches(0.4), Inches(0.4), [("Mye består", ROD, False)], size=15, font=HEAD, italic=True)
text(s, half + Inches(0.4), Inches(3.05), half - MX - Inches(0.4), Inches(1.2), "Faglig autoritet og varme", size=32, font=HEAD, bold=True)
text(s, half + Inches(0.4), Inches(4.1), half - MX - Inches(0.4), Inches(1.0), "Verdiene og menneskene er de samme.", size=17, font=BODY)
img(s, "logos/miles-logo-cream.png", MX, MY, width=Inches(1.0))

# 11 · Team
s = slide()
chrome(s)
text(s, Inches(0), Inches(1.5), EMU_W, Inches(0.5), [("Teamet ditt", ROD, False)], size=18, font=HEAD, italic=True, align=PP_ALIGN.CENTER)
members = [("Senior konsulent", "Navn Navnesen", "Løsningsarkitekt"),
           ("Konsulent", "Navn Navnesen", "Data scientist"),
           ("Senior konsulent", "Navn Navnesen", "UX-designer")]
d = Inches(1.7); gap = Inches(1.1); total = len(members) * d + (len(members) - 1) * gap
x0 = (EMU_W - total) / 2
for i, (role, nm, rl) in enumerate(members):
    cx = x0 + i * (d + gap)
    oval(s, cx, Inches(2.4), d, fill=KREM_DEEP, lineclr=BURG)
    pill(s, cx + d/2 - Inches(0.9), Inches(4.25), role, fill=ROD, fg=KREM, size=11, width=Inches(1.8))
    text(s, cx - Inches(0.4), Inches(4.8), d + Inches(0.8), Inches(0.4), [(nm, BURG, True)], size=15, font=BODY, align=PP_ALIGN.CENTER)
    text(s, cx - Inches(0.4), Inches(5.2), d + Inches(0.8), Inches(0.4), [(rl, RGBColor(0x6a,0x2f,0x44), False)], size=13, font=BODY, align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(6.3), EMU_W, Inches(0.4), [("Bytt de runde brønnene mot svart-hvitt-portretter.", RGBColor(0x6a,0x2f,0x44), False)], size=13, font=BODY, align=PP_ALIGN.CENTER, italic=True)

# ============================ DATA GRAPHICS ============================
# Native, editable PowerPoint charts (Edit Data still works) styled Tufte-minimal
# via tufte_pptx.style_chart. Illustrative data, mirroring charts.html.

def chart_slide(kick, title):
    s = slide()
    chrome(s)
    kicker(s, MX, MY + Inches(0.5), kick)
    text(s, MX, MY + Inches(0.92), Inches(11), Inches(0.9), title, size=30, font=HEAD, bold=True)
    return s

def cap(s, t):
    text(s, MX, Inches(6.7), Inches(11), Inches(0.4), [(t, tv.TINT2, False)], size=12, font=BODY, italic=True)

CRECT = (MX, Inches(2.0), EMU_W - 2 * MX, Inches(4.5))

# Divider
s = slide(BURG); chrome(s, dark=True)
rrect(s, Inches(2.6), Inches(2.7), Inches(8.1), Inches(2.1), fill=None, lineclr=KREM, radius=0.08)
text(s, Inches(2.6), Inches(2.95), Inches(8.1), Inches(0.4), [("Kapittel 03", KREM, False)], size=14, align=PP_ALIGN.CENTER)
text(s, Inches(2.6), Inches(3.45), Inches(8.1), Inches(1.2), [("Datagrafikk", KREM, False)], size=42, font=HEAD, bold=True, align=PP_ALIGN.CENTER)

# Nøkkeltall — sparkline-ish mini line charts
s = chart_slide("Nøkkeltall", "Siste 8 kvartaler")
kpi = [("Kundetilfredshet", [4.2,4.3,4.4,4.4,4.5,4.6,4.6,4.7], "4,7 / 5"),
       ("Ansatte totalt", [210,228,241,255,268,279,288,297], "297"),
       ("Beleggsgrad", [88,90,86,91,93,89,92,94], "94 %"),
       ("eNPS", [31,35,40,42,46,49,52,54], "54")]
for i, (nm, series, val) in enumerate(kpi):
    y = Inches(2.3 + i * 1.05)
    text(s, MX, y + Inches(0.15), Inches(4), Inches(0.5), [(nm, BURG, False)], size=20, font=BODY)
    tv.mini_line(s, (Inches(5.4), y, Inches(4.2), Inches(0.7)), series)
    text(s, Inches(10.0), y + Inches(0.12), Inches(2.5), Inches(0.5), [(val, BURG, True)], size=22, font=BODY, align=PP_ALIGN.RIGHT)
cap(s, "Sparklines som små native linjediagram. Illustrativ data.")

# Bar — revenue per service area
s = chart_slide("Sammenligning", "Omsetning per tjenesteområde")
tv.bar(s, CRECT,
       ["Transformasjon", "Brukeropplevelse", "Strategisk IT", "Data og AI", "Sky og plattform"],
       [88, 97, 142, 188, 206], accent_idx=3, unit=" MNOK")
cap(s, "Native stolpediagram fra null, verdien på stolpen, én rød aksent. Rediger data i PowerPoint. Illustrativ data.")

# Line — employee growth
s = chart_slide("Utvikling over tid", "Ansatte per tjenesteområde")
lc = tv.line(s, CRECT, ["2019","2020","2021","2022","2023","2024","2025","2026"],
             [("Data og AI", [18,26,38,55,74,96,121,148]),
              ("Sky og plattform", [64,71,78,84,90,97,104,110]),
              ("Strategisk IT", [41,44,46,48,49,51,52,54])],
             accent_name="Data og AI")
try:
    lc.value_axis.minimum_scale = 0; lc.value_axis.major_unit = 50
except Exception:
    pass
cap(s, "Native linjediagram, accent-serien i rødt, resten trukket tilbake. Illustrativ data.")

# Clustered bar — consultants per office, before/after
s = chart_slide("Rangering", "Konsulenter per kontor (2024 → 2025)")
tv.clustered_bar(s, CRECT,
                 ["Innlandet","Haugesund","Ålesund","Litauen","Stavanger","Trondheim","Bergen","Oslo"],
                 [("2024", [6,9,11,14,31,30,54,112]), ("2025", [8,9,12,26,29,34,58,121])])
cap(s, "To serier side om side (før/etter), minimal forklaring. Illustrativ data.")

# Slopegraph — AI adoption (2-category line)
s = chart_slide("Før og etter", "Andel prosjekter med AI-komponent (%)")
tv.slope(s, (Inches(3.0), Inches(2.0), Inches(7.3), Inches(4.5)),
         [("Data og AI", 62, 91), ("Sky og plattform", 28, 54), ("Strategisk IT", 12, 38),
          ("Brukeropplevelse", 9, 31), ("Transformasjon", 5, 22)],
         "2023", "2025", accent_name="Data og AI", value_fmt='0"%"')
cap(s, "Skråningsgraf = 2-kategori linjediagram, én serie per område. Illustrativ data.")

# Scatter
s = chart_slide("Korrelasjon", "Varighet vs. teamstørrelse")
tv.scatter(s, CRECT,
           [(2,6),(3,7),(3,10),(4,9),(4,14),(5,12),(5,16),(6,13),(6,20),(7,18),(7,24),
            (8,19),(8,28),(9,26),(10,31),(11,38),(12,34),(4,8),(9,22)],
           highlight=(5, 11), x_label="Teamstørrelse", y_label="Varighet (uker)")
cap(s, "Native XY-punktdiagram, uthevet referansepunkt i rødt. Illustrativ data.")

# Small multiples — grid of mini line charts
s = chart_slide("Mange serier", "Vekst per kontor 2021–2026")
sm = [("Oslo",[98,104,109,112,117,121]),("Bergen",[44,47,50,54,56,58]),
      ("Trondheim",[22,25,27,30,32,34]),("Stavanger",[27,28,30,31,30,29]),
      ("Litauen",[4,7,10,14,20,26]),("Ålesund",[8,9,10,11,12,12]),
      ("Haugesund",[6,7,8,9,9,9]),("Innlandet",[3,4,5,6,7,8])]
PW, PH, GX, GY = Inches(2.7), Inches(1.45), Inches(0.25), Inches(0.55)
x0, y0 = MX, Inches(2.3)
for i, (nm, vals) in enumerate(sm):
    col, row = i % 4, i // 4
    px = Emu(int(x0) + col * (int(PW) + int(GX)))
    py = Emu(int(y0) + row * (int(PH) + int(GY) + int(Inches(0.35))))
    text(s, px, py, PW, Inches(0.3), [(nm, BURG, True)], size=13, font=BODY)
    tv.mini_line(s, (px, Emu(int(py) + int(Inches(0.32))), PW, PH), vals, y_range=(0, 125))
cap(s, "Small multiples: åtte native mini-linjediagram på felles skala (0–125). Illustrativ data.")

# 12 · Closing red
s = slide(ROD)
img(s, "logos/miles-logo-cream.png", MX, MY, width=Inches(1.6))
text(s, MX, Inches(3.9), Inches(10), Inches(1.2), [("Takk for tiden din! :)", KREM, False)], size=48, font=HEAD, bold=True)
labels = ["←  Forside", "←  Innholdsoversikten", "←  Har du noe på hjertet?"]
x = MX
for lb in labels:
    p = pill(s, x, Inches(5.4), lb, fill=KREM, fg=BURG, size=13, width=Inches(0.11) * len(lb) + Inches(0.5))
    x += p.width + Inches(0.25)

out = os.path.join(HERE, "miles-templates.pptx")
prs.save(out)
print("wrote", out, "—", len(prs.slides._sldIdLst), "slides")
