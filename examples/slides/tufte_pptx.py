"""
tufte_pptx.py — make native, editable PowerPoint charts that look Tufte-minimal
and Miles-on-brand.

The whole point: these are REAL PowerPoint charts (python-pptx `add_chart`), so
"Edit Data" still opens the embedded spreadsheet and the chart re-renders — they
behave like any PowerPoint chart, they just look better. We don't draw pictures.

`style_chart()` is the reusable pass — call it on ANY python-pptx chart to get
the house look: no gridlines, no borders, no chart/plot-area fill, no tick
marks, hairline (or absent) axis lines, DM Sans labels in burgundy, values
printed on the marks, and exactly one Miles-red accent via per-point/series fill.

Two Tufte refinements can't survive a user data-edit in native PowerPoint —
true range-frames (axis spanning only data min..max with labels there) and
direct line-end labels. We trade those for editability and lean on a minimal,
border-less legend only where multi-series naming is unavoidable.
"""
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK,
                             XL_LEGEND_POSITION, XL_MARKER_STYLE)
from pptx.oxml.ns import qn

# Miles palette
KREM = RGBColor(0xFB, 0xF0, 0xE5)
BURG = RGBColor(0x45, 0x0D, 0x21)
ROD = RGBColor(0xFF, 0x30, 0x3B)
TINT1 = RGBColor(0x6A, 0x2F, 0x44)
TINT2 = RGBColor(0x8F, 0x52, 0x65)
TINT3 = RGBColor(0xB5, 0x85, 0x94)
BODY = "DM Sans"


def _nofill(el, tagns='c'):
    """Give a chartSpace/plotArea element a transparent fill + no border, so the
    cream slide shows through (no white box, no frame). Schema-position aware."""
    try:
        spPr = el.find(qn(f'{tagns}:spPr'))
        if spPr is None:
            spPr = el.makeelement(qn(f'{tagns}:spPr'), {})
            if tagns == 'c' and el.tag == qn('c:chartSpace'):
                anchor = el.find(qn('c:chart'))
                anchor.addnext(spPr)          # spPr follows <c:chart>
            else:
                el.append(spPr)               # plotArea: spPr is last
        spPr.append(spPr.makeelement(qn('a:noFill'), {}))
        ln = spPr.makeelement(qn('a:ln'), {})
        ln.append(ln.makeelement(qn('a:noFill'), {}))
        spPr.append(ln)
    except Exception:
        pass  # degrade gracefully: chart just keeps its default white area


def _line_nofill(fmt):
    try:
        fmt.line.fill.background()
    except Exception:
        pass


def style_chart(chart, *, value_fmt=None, labels=False, label_pos=None,
                show_value_axis=False, show_cat_axis=True, legend=False,
                font_size=11):
    """The reusable Tufte-Miles pass. Erase chartjunk; keep data + labels."""
    chart.has_title = False
    chart.font.name = BODY
    chart.font.size = Pt(font_size)
    chart.font.color.rgb = BURG

    chart.has_legend = bool(legend)
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_size)
        chart.legend.font.color.rgb = BURG

    for axis, visible in ((chart.value_axis, show_value_axis),
                          (chart.category_axis, show_cat_axis)):
        try:
            axis.has_major_gridlines = False
            axis.has_minor_gridlines = False
            axis.major_tick_mark = XL_TICK_MARK.NONE
            axis.minor_tick_mark = XL_TICK_MARK.NONE
            axis.visible = bool(visible)
            _line_nofill(axis.format)
            axis.tick_labels.font.size = Pt(font_size)
            axis.tick_labels.font.name = BODY
            axis.tick_labels.font.color.rgb = TINT2
        except Exception:
            pass

    if labels:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        if value_fmt:
            dl.number_format = value_fmt
            dl.number_format_is_linked = False
        if label_pos is not None:
            try: dl.position = label_pos
            except Exception: pass
        dl.font.size = Pt(font_size)
        dl.font.name = BODY
        dl.font.color.rgb = BURG

    # erase the white chart-area box and the plot-area frame
    _nofill(chart._chartSpace)
    pa = chart._chartSpace.find(qn('c:chart')).find(qn('c:plotArea'))
    if pa is not None:
        _nofill(pa)
    return chart


# ----------------------------------------------------------------- chart builders
def bar(slide, rect, cats, vals, *, accent_idx=None, unit='', value_fmt='0'):
    """Horizontal bar from zero, value printed on each bar, one accent bar."""
    cd = CategoryChartData(); cd.categories = cats
    cd.add_series('verdi', vals, number_format=f'0"{unit}"' if unit else '0')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, *rect, cd)
    ch = gf.chart
    s = ch.series[0]
    s.format.fill.solid(); s.format.fill.fore_color.rgb = BURG
    s.format.line.fill.background()
    if accent_idx is not None:
        pt = s.points[accent_idx]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb = ROD
    ch.plots[0].gap_width = 80
    style_chart(ch, labels=True, label_pos=XL_LABEL_POSITION.OUTSIDE_END,
                value_fmt=(f'0"{unit}"' if unit else '0'), show_value_axis=False,
                show_cat_axis=True)
    return ch


def clustered_bar(slide, rect, cats, series, *, accent_name=None):
    """Two series side by side (e.g. before/after). Minimal legend names them."""
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals, number_format='0')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, *rect, cd)
    ch = gf.chart
    palette = [TINT3, BURG]  # earlier series faint, latest series strong
    for i, s in enumerate(ch.series):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = palette[i % len(palette)]
        s.format.line.fill.background()
    ch.plots[0].gap_width = 60
    ch.plots[0].overlap = -10
    style_chart(ch, labels=False, show_value_axis=False, show_cat_axis=True, legend=True)
    return ch


def line(slide, rect, cats, series, *, accent_name=None, value_fmt='0'):
    """Multi-series time line. Accent series in red, the rest recede; minimal legend."""
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals, number_format='0')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, *rect, cd)
    ch = gf.chart
    recede = [BURG, TINT2, TINT3]
    j = 0
    for s in ch.series:
        s.smooth = False
        s.marker.style = XL_MARKER_STYLE.NONE
        if s.name == accent_name:
            s.format.line.color.rgb = ROD; s.format.line.width = Pt(2.5)
        else:
            s.format.line.color.rgb = recede[j % len(recede)]; s.format.line.width = Pt(1.5); j += 1
    style_chart(ch, labels=False, show_value_axis=True, show_cat_axis=True, legend=True)
    return ch


def slope(slide, rect, items, left_label, right_label, *, accent_name=None, value_fmt='0'):
    """Slopegraph = a 2-category line chart, one series per item. Values at both ends."""
    cd = CategoryChartData(); cd.categories = [left_label, right_label]
    for name, l, r in items:
        cd.add_series(name, (l, r), number_format='0')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, *rect, cd)
    ch = gf.chart
    for s in ch.series:
        s.smooth = False
        s.marker.style = XL_MARKER_STYLE.CIRCLE; s.marker.size = 5
        col = ROD if s.name == accent_name else TINT2
        s.format.line.color.rgb = col; s.format.line.width = Pt(2)
        s.marker.format.fill.solid(); s.marker.format.fill.fore_color.rgb = col
        s.marker.format.line.fill.background()
    style_chart(ch, labels=True, value_fmt=value_fmt, show_value_axis=False,
                show_cat_axis=True, legend=True)
    return ch


def scatter(slide, rect, pts, *, highlight=None, x_label='', y_label=''):
    """XY scatter, markers only. Highlight point as a second one-point red series."""
    xy = XyChartData()
    main = xy.add_series('punkter', number_format='0')
    for x, y in pts:
        main.add_data_point(x, y)
    if highlight is not None:
        hi = xy.add_series('uthevet', number_format='0')
        hi.add_data_point(highlight[0], highlight[1])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, *rect, xy)
    ch = gf.chart
    ser = ch.series
    ser[0].marker.style = XL_MARKER_STYLE.CIRCLE; ser[0].marker.size = 6
    ser[0].marker.format.fill.solid(); ser[0].marker.format.fill.fore_color.rgb = BURG
    ser[0].marker.format.line.fill.background()
    ser[0].format.line.fill.background()
    if highlight is not None:
        ser[1].marker.style = XL_MARKER_STYLE.CIRCLE; ser[1].marker.size = 11
        ser[1].marker.format.fill.solid(); ser[1].marker.format.fill.fore_color.rgb = ROD
        ser[1].marker.format.line.color.rgb = ROD
        ser[1].format.line.fill.background()
    style_chart(ch, labels=False, show_value_axis=True, show_cat_axis=True, legend=False)
    try:
        ch.category_axis.axis_title.text_frame.text = x_label
        ch.value_axis.axis_title.text_frame.text = y_label
        for ttl in (ch.category_axis.axis_title, ch.value_axis.axis_title):
            f = ttl.text_frame.paragraphs[0].runs[0].font
            f.name = BODY; f.size = Pt(11); f.bold = False; f.color.rgb = TINT2
    except Exception:
        pass
    return ch


def mini_line(slide, rect, vals, *, y_range=None):
    """A sparkline-ish tiny line: no axes, no legend, last point marked red.
    Pass y_range=(lo,hi) to put several minis on a SHARED scale — required for
    honest small-multiples (magnitude comparable across panels). Omit it for
    standalone sparklines, which show shape and auto-fit by design."""
    cd = CategoryChartData(); cd.categories = [str(i) for i in range(len(vals))]
    cd.add_series('s', vals, number_format='0')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, *rect, cd)
    ch = gf.chart
    s = ch.series[0]
    s.smooth = False; s.marker.style = XL_MARKER_STYLE.NONE
    s.format.line.color.rgb = BURG; s.format.line.width = Pt(1.75)
    if y_range is not None:
        try:
            ch.value_axis.minimum_scale = y_range[0]
            ch.value_axis.maximum_scale = y_range[1]
        except Exception:
            pass
    pt = s.points[len(vals) - 1]
    try:
        pt.marker.style = XL_MARKER_STYLE.CIRCLE; pt.marker.size = 6
        pt.marker.format.fill.solid(); pt.marker.format.fill.fore_color.rgb = ROD
    except Exception:
        pass
    style_chart(ch, labels=False, show_value_axis=False, show_cat_axis=False, legend=False)
    return ch
